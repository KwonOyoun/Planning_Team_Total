# main.py
import os
import html
import re
import requests
from pptx.util import Pt
from pptx import Presentation
from xml.etree import ElementTree as ET

from .config import HWPX_PATH, PPTX_PATH, OUTPUT_DIR, GOOGLE_API_KEY, GOOGLE_CX

NS = {"hp": "http://www.hancom.co.kr/hwpml/2011/paragraph"}
HP = "{http://www.hancom.co.kr/hwpml/2011/paragraph}"


from .parser import unzip_hwpx, load_xml, save_xml, zip_hwpx, parse_need, parse_overview
from .paragraph_editor import (
    find_insert_index,
    find_insert_index_2,
    find_insert_index_3,
    find_insert_index_4,
    insert_need_paragraphs,
    insert_ref_paragraphs,
    modify_text_contents,
    modify_text_contents_ref,
    find_insert_index_3_1, #[목차]
    # apply_font_to_business_section
)
from .web_researcher import get_fact_sheet
from .biz_generator import (
    generate_titles,
    generate_overview,
    generate_need,
    generate_suggestion,
    generate_reference1,
    generate_reference2,
    generate_reference3,
    generate_reference4,
    summarize_business_evidences    
)


def extract_items_ref(text: str):
    # 사업개요 추출
    m1 = re.search(r"사업개요\s*:\s*(.*)", text)
    a = m1.group(1).strip() if m1 else ""

    # 수혜자 추출
    m2 = re.search(r"수혜자\s*:\s*(.*)", text)
    b = m2.group(1).strip() if m2 else ""

    return a, b



def extract_items(text: str):
    """
    '1. 내용' 형식으로 된 항목들을 번호 제거하고
    내용만 리스트로 반환하는 함수
    """
    # 줄 단위로 split
    lines = text.strip().splitlines()

    items = []
    for line in lines:
        # 앞의 번호 제거 (숫자. 또는 숫자) 같은 패턴)
        cleaned = re.sub(r'^\s*\d+\.\s*', '', line).strip()
        if cleaned:
            items.append(cleaned)

    return items

# from parser import unzip_hwpx, load_xml, save_xml, zip_hwpx
# from paragraph_editor import find_insert_index, insert_need_paragraphs
# from config import HWPX_PATH, OUTPUT_PATH

# ------------------------------------------------------------------------
# 장표 제작 함수들
# ------------------------------------------------------------------------
# 1 재귀적으로 모든 shape 순회 (GroupShape 포함)
def iter_all_shapes(shapes):
    for shape in shapes:
        yield shape
        if shape.shape_type == 6:  # GROUP shape
            for nested in iter_all_shapes(shape.shapes):
                yield nested

# 2. 텍스트 템플릿에 맞게 후처리 
def parse_business_text(text, title, purpose):
    """
    입력된 장문의 사업 설명 텍스트(ref3 형식)를 분석하여
    사업내용 / 사업주장 / 사업근거 블록을 자동 추출하여 dict 형태로 반환

    - 사업근거는 섹션당 최대 2개 유지
    - 각 근거는 40자 이내로 OpenAI 요약
    - OpenAI 호출은 섹션당 1회 (토큰 최소화)
    """

    # 줄 단위 정리 (빈 줄 제거)
    lines = [line.strip() for line in text.split("\n") if line.strip()]

    replacements = {
        "(사업명)": title,
        "(사업목표)": purpose
    }

    section_index = 0
    current_근거 = []

    for raw in lines:
        line = raw.strip()

        # -------------------------------------------------
        # 1) "ㅇ (A) B" → 사업내용 / 사업주장
        # -------------------------------------------------
        m = re.match(r"^ㅇ\s*\((.*?)\)\s*(.*)$", line)
        if m:
            # 🔥 이전 섹션의 근거 처리
            if current_근거:
                summarized = summarize_business_evidences(
                    current_근거,
                    max_len=40
                )

                # 최대 2개 유지
                summarized = summarized[:2]

                replacements[f"(사업근거{section_index})"] = "\n\n".join(
                    [f"{i+1}. {s}" for i, s in enumerate(summarized)]
                )

                current_근거 = []

            # 새 섹션 시작
            section_index += 1
            내용 = m.group(1).strip()
            주장 = m.group(2).strip()

            replacements[f"(사업내용{section_index})"] = 내용
            replacements[f"(사업주장{section_index})"] = 주장
            continue

        # -------------------------------------------------
        # 2) "- (C) D" → 사업근거 수집
        # -------------------------------------------------
        m = re.match(r"^-+\s*\((.*?)\)\s*(.*)$", line)
        if m:
            근거문장 = f"{m.group(1).strip()} {m.group(2).strip()}"

            # 🔥 최대 2개까지만 수집
            if len(current_근거) < 2:
                current_근거.append(근거문장)

            continue

    # -------------------------------------------------
    # 🔥 마지막 섹션 근거 처리 (중요)
    # -------------------------------------------------
    if current_근거:
        summarized = summarize_business_evidences(
            current_근거,
            max_len=40
        )

        summarized = summarized[:2]

        replacements[f"(사업근거{section_index})"] = "\n\n".join(
            [f"{i+1}. {s}" for i, s in enumerate(summarized)]
        )

    return replacements


# 3. 대체 및 스타일 적용
def replace_and_style_pptx(
    input_path,
    output_path,
    replacements,
    evidence_font_name=None,
    evidence_font_size=None
):
    prs = Presentation(input_path)

    for slide in prs.slides:
        for shape in iter_all_shapes(slide.shapes):

            # ------------------ 텍스트 상자 ------------------
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:

                    original = p.text
                    new = original
                    
                    is_evidence = False
                    is_content = False

                    # 치환
                    for key, value in replacements.items():
                        key_norm = key.replace(" ", "")
                        text_norm = new.replace(" ", "")

                        if key_norm in text_norm:
                            new = new.replace(key, value)
                            new = new.replace(key.replace(" ", ""), value)

                            # 사업근거이면 표시
                            if "사업근거" in key_norm:
                                is_evidence = True

                            if "사업주장" in key_norm:
                                is_evidence = True

                            if "사업내용" in key_norm:
                                is_content = True

                    # 변경된 경우 run 재구성 + 스타일 적용
                    if new != original:
                        # 기존 run 삭제
                        for _ in range(len(p.runs)):
                            p._p.remove(p.runs[0]._r)

                        # 새 run 생성
                        run = p.add_run()
                        run.text = new

                        # 만약 사업근거 문단이면 스타일 적용
                        if is_evidence:
                            run.font.size = Pt(12)
                        if is_content:
                            run.font.size = Pt(14)


            # ------------------ 표 ------------------
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        for p in cell.text_frame.paragraphs:
                            
                            original = p.text
                            new = original

                            is_evidence = False
                            is_content = False

                            for key, value in replacements.items():
                                key_norm = key.replace(" ", "")
                                text_norm = new.replace(" ", "")

                                if key_norm in text_norm:
                                    new = new.replace(key, value)
                                    new = new.replace(key.replace(" ", ""), value)

                                    # 사업근거이면 표시
                                    if "사업근거" in key_norm:
                                        is_evidence = True

                                    if "사업주장" in key_norm:
                                        is_evidence = True

                                    if "사업내용" in key_norm:
                                        is_content = True

                            if new != original:
                                for _ in range(len(p.runs)):
                                    p._p.remove(p.runs[0]._r)

                                run = p.add_run()
                                run.text = new

                                if is_evidence:
                                    run.font.size = Pt(12)
                                if is_content:
                                    run.font.size = Pt(14)

    prs.save(output_path)

#####################################################################
#####################################################################



# ------------------------------------------------------------------------
# 목차 제작 함수들
# ------------------------------------------------------------------------
def make_p(text: str):
    NS_HP = "{%s}" % NS['hp']

    # 문단(p) 생성
    p = ET.Element(f"{NS_HP}p")
    p.set("styleIDRef", "0")     # ← style 적용
    p.set("paraPrIDRef", "0")    # ← paragraph 속성 적용
    p.set("pageBreak", "0")
    p.set("columnBreak", "0")
    p.set("merged", "0")

    # run 생성
    run = ET.SubElement(p, f"{NS_HP}run")
    run.set("charPrIDRef", "51")   # ← charPrID 적용

    # 텍스트 추가
    t = ET.SubElement(run, f"{NS_HP}t")
    t.text = text

    return p

def convert_ref3_to_numbered(ref3: str) -> list:
    """
    ref3을 배열(list) 형태의 번호문단 리스트로 변환
    ※ 목차용: 3.X 까지만 생성
    ※ 괄호 뒤 설명을 우선 사용
    """
    lines = [l.rstrip() for l in ref3.split("\n") if l.strip()]

    result = []
    top_idx = 0

    for line in lines:

        # ㅇ (괄호) 뒤 설명 캡처
        m = re.match(r"^ㅇ\s*\((.*?)\)\s*(.*)$", line)
        if m:
            top_idx += 1
            bracket = m.group(1).strip()
            after = m.group(2).strip()

            # 🔥 괄호 뒤 설명 우선 사용
            title = after if after else bracket

            result.append(f"  3.{top_idx} {title}")

    return result



def replace_business_content(root, ref3):
    """
    HWPX 문서에서 [사업내용] 문단 위치에
    번호형 문단을 여러 개 생성해서 삽입
    """

    numbered_lines = convert_ref3_to_numbered(ref3)

    # ① [사업내용] 문단 찾기
    target_p = None
    p_list = list(root)

    for p in p_list:
        t = p.find(".//hp:t", NS)
        if t is not None and "[사업내용]" in t.text:
            target_p = p
            break

    if target_p is None:
        print("⚠ [사업내용] 위치를 찾을 수 없습니다.")
        return

    # target 위치
    insert_pos = p_list.index(target_p)

    # ② 기존 문단 삭제
    root.remove(target_p)

    # ③ 번호문단 여러 개 삽입
    for i, line in enumerate(numbered_lines):
        new_p = make_p(line)
        root.insert(insert_pos + i, new_p)
#####################################################################
#####################################################################


def generate_titles_from_keywords(keywords_string: str, prompts: dict) -> list:
    """
    프론트엔드에서 넘어온 키워드 배열을 받아
    GPT에게 제목 생성을 요청하고, 결과 텍스트를 그대로 반환.
    """
    titles_text_ = generate_titles(keywords_string, prompts["title"])  # 기존에 있던 함수 그대로 사용
    titles_text = extract_items(titles_text_) # 배열로 가공
    return titles_text

# html을 보기 좋게 출력하기 위한 함수
def ref_to_html_formatted(ref_text):
    lines = [l.rstrip() for l in ref_text.split("\n") if l.strip()]
    html_lines = []

    started = False  # 첫 번째 'ㅇ'에서는 줄 띄우기 금지

    for line in lines:
        clean = line.lstrip()

        match = re.match(r"^(ㅇ|-|\*)\s*(.*)", clean)
        if not match:
            continue

        symbol = match.group(1)
        content = html.escape(match.group(2).strip())

        # 들여쓰기 규칙
        if symbol == "ㅇ":
            indent = ""
        elif symbol == "-":
            indent = "&nbsp;&nbsp;"
        elif symbol == "*":
            indent = "&nbsp;&nbsp;&nbsp;&nbsp;"
        else:
            indent = ""

        # 🔥 핵심 수정점: 'ㅇ'이 두 번째 이후 등장하면 줄 1줄만 띄우기
        if symbol == "ㅇ" and started:
            html_lines.append("")  # 빈 줄 하나 → 나중에 join에서 <br> 처리됨

        # 라인 추가
        html_lines.append(f"{indent}{symbol} {content}")

        if symbol == "ㅇ":
            started = True

    # 🔥 join 방식 변경 (한 줄을 <br>로만 바꿈)
    return "<br>".join(html_lines)




# 웹에서 호출할 함수 (키워드 배열 → 제목 후보들)
def generate_proposal_and_hwpx(title: str, keywords_string: str, prompts: dict) -> str:
    """
    프론트엔드에서 넘어온 키워드 배열을 받아
    GPT에게 제목 생성을 요청하고, 결과 텍스트를 그대로 반환.
    """

    print("=== 제안서 생성 요청 중 ===")
    print("선택된 제목:", title)
    print("키워드 문자열:", keywords_string)

    # 1) 개요 생성
    overview = generate_overview(title, prompts["overview"])
    print(overview)

    # 1-1) 개요 파싱 (사업 목적 / 사업 내용)
    purpose = ""
    content = ""
    for line in overview.split("\n"):
        if "■ 사업 목적" in line:
            continue
        if "■ 사업 내용" in line:
            continue

        if line.startswith("ㅇ") and not content:
            purpose += line + "\n"
        elif line.startswith("ㅇ"):
            content += line + "\n"

    purpose, content = parse_overview(overview)


    # 2) 필요성 생성
    need = generate_need(title, purpose, content, prompts["need"])
    print(need)

    need_parts = parse_need(need)


    # 3) 건의 사항 생성
    suggestion = generate_suggestion(title, purpose, content, need, prompts["suggestion"])
    print("\n■ 건의 문장")
    print(suggestion)

    # # 테스트용 더미 데이터
    # title = title
    # overview= 'content\npurpose'
    # content = 'content'
    # purpose = 'purpose'
    # need = '[1) 필요성]\nㅇ AI 기반 의료기기 개발의 필요성은 분석 자동화 수요 때문입니다.\n- 국내 병원의 의료데이터 분석 인력 부족이 주요 원인입니다.\n* 2024년 기준 의료 AI 시장은 연 17% 성장하고 있습니다.\n[2) 필요성]\nㅇ 다기관 데이터 활용 필요성이 증가하고 있습니다.\n- 병원 간 데이터 편차가 모델 성능 저하를 초래합니다.\n* FDA도 다기관 기반 검증을 요구하고 있습니다.\n[3) 필요성]\nㅇ 디지털헬스케어 규제 완화는 혁신기술 수요 증가와 연계됩니다.\n- 기업들은 인허가 대응 비용이 증가하고 있습니다.\n* 정부는 2025년 디지털 규제 혁신 43건을 발표했습니다.'
    # suggestion = 'suggestion'

    
    # ------------------------------------------------------------------------
    # 2. 참고 자료 부분 작성
    # ------------------------------------------------------------------------
    print("\n----------------- 참고 자료 -----------------")

    print("\n1. 사업 개요")
    ref1 = generate_reference1(title, purpose, content, prompts["ref1"])
    print(ref1)

    print("\n2. 추진 배경 및 필요성 (RAG 수행 중...)")
    # RAG: 검색 수행
    context_ref2 = get_fact_sheet(f"{title} 시장 규모 통계 필요성")
    ref2 = generate_reference2(title, purpose, content, need, prompts["ref2"])
    print(ref2)

    print("\n3. 사업 내용 (RAG 수행 중...)")
    # RAG: 검색 수행
    context_ref3 = get_fact_sheet(f"{title} 기술 동향 구축 사례")
    ref3 = generate_reference3(title, purpose, content, need, prompts["ref3"])
    print(ref3)

    print("\n4. 기대 효과 (RAG 수행 중...)")
    # RAG: 검색 수행
    context_ref4 = get_fact_sheet(f"{title} 경제적 사회적 기대효과")
    ref4 = generate_reference4(title, purpose, content, need, prompts["ref4"])
    print(ref4)


    # ------------------------------------------------------------------------
    # 4. HWPX 파일 생성 (GPT 재호출 없이, 위 텍스트 사용)
    # ------------------------------------------------------------------------

    over_ref, benef = extract_items_ref(ref1)

    # 필요성 문단 분리
    need_parts = parse_need(need)

    # 템플릿 언팩 & XML 로딩
    unzip_hwpx(HWPX_PATH)
    tree, root = load_xml()
    p_list = list(root)

    # '필요성' 삽입 위치 찾기 & 삽입
    insert_index = find_insert_index(p_list)
    insert_need_paragraphs(root, insert_index, need_parts)

    p_list = list(root)

    # '참고' 삽입 위치 찾기 & 삽입
    insert_index_2 = find_insert_index_2(p_list) # 2
    source = insert_ref_paragraphs(root, insert_index_2, ref2) # 2

    p_list = list(root)
    
    insert_index_3 = find_insert_index_3(p_list) # 3
    source2 = insert_ref_paragraphs(root, insert_index_3, ref3) # 3

    p_list = list(root)

    insert_index_4 = find_insert_index_4(p_list) # 4
    source3 = insert_ref_paragraphs(root, insert_index_4, ref4) # 4
    

    # 제목/개요/필요성/제안내용 텍스트 반영
    modify_text_contents(
        root,
        title=title,
        purpose=purpose,
        content=content,
        suggestion=suggestion
    )

    modify_text_contents_ref(
        root,
        title = title,
        overview = over_ref,
        benef = benef
    )

    # 목차 제작 -> 사업내용 치환
    # insert_index3_1 = find_insert_index_3_1(p_list) # 4
    # insert_ref_paragraphs(root,insert_index3_1, )
    replace_business_content(root, ref3)


    # ------------------------------------------------------------------------
    # 4. '삭제요망' 태그가 붙은 문단 삭제(문단 스타일 저장을 위해 남겨둔 것들)
    # ------------------------------------------------------------------------

    NS = {"hp": "http://www.hancom.co.kr/hwpml/2011/paragraph",}
    delete_targets = []

    for p in root.findall(".//hp:p", NS):
        t = p.find(".//hp:t", NS)
        if t is not None and "삭제요망" in t.text:
            delete_targets.append(p)

    # 부모에서 삭제하기
    for p in delete_targets:
        parent = p.getparent() if hasattr(p, "getparent") else root
        # ElementTree는 getparent() 기본 제공 안하니까 아래 방식 사용
        for parent in root.iter():
            if p in list(parent):
                parent.remove(p)
                break

    # ------------------------------------------------------------------------
    # 5. PPTX로 장표 제작
    # ------------------------------------------------------------------------
    # 따로 내용에 추가
    replacements = parse_business_text(ref3, title, purpose)

    os.makedirs(OUTPUT_DIR, exist_ok=True)

    replace_and_style_pptx(
        input_path=PPTX_PATH,
        output_path=os.path.join(OUTPUT_DIR, f"{title}.pptx"),
        replacements=replacements,
        evidence_font_name="맑은 고딕",
        evidence_font_size=11
    )

    def sanitize_filename(name: str) -> str:
        return re.sub(r'[\\/:*?"<>|]', '', name).strip()
    
    # 저장 & 다시 압축
    save_xml(tree)
    safe_title = sanitize_filename(title)
    zip_hwpx(os.path.join(OUTPUT_DIR, f"{safe_title}.hwpx"))

    print("✔ 제안서 생성 + HWPX 반영 완료:", f"output/{title}.hwpx")


    # ------------------------------------------------------------------------
    # 6. 출처(References) 별도 파일 저장 (구글 검색 + 인용 문구 찾기)
    # ------------------------------------------------------------------------
        
    # 실제 링크 검색 함수 (Google Custom Search API 사용)
    def find_real_url(query):
        # API 키가 없으면 검색 불가
        if not GOOGLE_API_KEY or not GOOGLE_CX:
            # print("⚠ Google API Key 또는 CX가 설정되지 않아 검색을 건너뜁니다.")
            return None

        try:
            # 출처 텍스트에서 불필요한 기호 제거하고 검색
            clean_query = re.sub(r'https?://\S+', '', query) # 기존 URL 제거
            clean_query = re.sub(r'[^\w\s가-힣]', '', clean_query) # 특수문자 제거
            clean_query = clean_query.strip()
            
            if not clean_query:
                return None
            
            # 너무 짧으면 검색 스킵
            if len(clean_query) < 2:
                return None

            print(f"🔎 Google API 검색 중: {clean_query[:20]}...")
            
            # Google Custom Search API 호출
            url = "https://www.googleapis.com/customsearch/v1"
            params = {
                'key': GOOGLE_API_KEY,
                'cx': GOOGLE_CX,
                'q': clean_query,
                'num': 1  # 1개만 가져옴
            }
            
            response = requests.get(url, params=params, timeout=5)
            
            if response.status_code == 200:
                data = response.json()
                if 'items' in data and len(data['items']) > 0:
                    found_link = data['items'][0]['link']
                    print(f"   → 찾은 링크: {found_link}")
                    return found_link
                else:
                    print("   → 검색 결과 없음")
                    return None
            else:
                print(f"⚠ Google API 오류: {response.status_code} {response.text}")
                return None

        except Exception as e:
            print(f"⚠ 검색 로직 예외 발생: {e}")
            return None
        return None

    # 본문에서 [1. 키워드] 또는 [1] 형태의 인용구를 찾아 문장을 추출하는 함수
    def get_citation_context(full_text):
        context_map = {} # { 1: ["문장1", "문장2"], 2: ... }
        
        lines = full_text.split('\n')
        for line in lines:
            line = line.strip()
            if not line: continue
            
            # [숫자. 형태 찾기
            matches = re.findall(r'\[(\d+)\.', line)
            for m in matches:
                idx = int(m)
                if idx not in context_map:
                    context_map[idx] = []
                context_map[idx].append(line)
        
        return context_map

    # 섹션별 처리 함수
    def process_section_references(f, section_name, ref_list, full_text):
        if not ref_list:
            return

        f.write(f"\n[{section_name} 참고자료]\n")
        
        # 본문에서 인용된 위치 찾기
        context_map = get_citation_context(full_text)

        for idx, r in enumerate(ref_list, 1):
            # 1. 텍스트 추출
            text_part = re.sub(r'https?://\S+', '', r).strip()
            
            # 2. 검색어 생성
            search_query = text_part
            if not text_part:
                    search_query = f"{title} 관련 자료 통계 시장규모"

            # 3. 실제 검색
            new_url = find_real_url(search_query)
            
            # 4. 파일 쓰기
            f.write(f"{idx}. 항목: {text_part if text_part else '(자동생성 키워드: '+search_query+')'}\n")
            
            # 인용된 문장 표시
            # ref_list의 인덱스(1,2,3...)와 본문의 [1], [2]... 매칭
            # 보통 순서대로 나오므로 idx를 사용
            if idx in context_map:
                f.write(f"   ▶ 인용된 내용:\n")
                for ctx in context_map[idx]:
                        f.write(f"     - \"{ctx}\"\n")
            else:
                f.write(f"   ▶ (본문에서 인용 마커 [{idx}.]를 찾지 못함)\n")

            if new_url:
                f.write(f"   └─ [검색된 링크] {new_url}\n")
            else:
                f.write(f"   └─ [링크 찾기 실패]\n")
            f.write("-" * 50 + "\n")


    # 파일로 저장 (섹션별로 구분)
    ref_file_path = f"output/{title}_references.txt"
    with open(ref_file_path, "w", encoding="utf-8") as f:
        f.write(f"=== {title} 참고 자료 출처 (AI 검색 + 인용 분석) ===\n")
        f.write("※ 주의: 아래 '인용된 내용'은 AI가 생성한 문장이며, '검색된 링크'는 해당 내용을 검증하기 위해 별도로 검색한 결과입니다.\n\n")
        
        process_section_references(f, "2. 추진 배경 및 필요성", source, ref2)
        process_section_references(f, "3. 사업 내용", source2, ref3)
        process_section_references(f, "4. 기대 효과", source3, ref4)

    print("✔ 참고 자료 출처 저장 완료 (구글 검색 적용):", ref_file_path)


    # --------------------------------------------------------
    # HWPX 파일에 참고자료 반영
    # --------------------------------------------------------
    def parse_reference_txt(txt: str):
        """
        참고자료 TXT를 파싱하여
        [
        {
            "no": 1,
            "title": "...",
            "quotes": [...],
            "link": "..."
        },
        ...
        ]
        형태로 반환
        """
        results = []

        blocks = re.split(r"\n\s*\d+\.\s*항목:", txt)
        blocks = blocks[1:]  # 헤더 제거

        for idx, block in enumerate(blocks, start=1):
            title_match = re.search(r"\[(.*?)\]", block)
            if not title_match:
                continue
            title = title_match.group(1).strip()

            quotes = re.findall(r'-\s*"([^"]+)"', block)

            link_match = re.search(r"\[검색된 링크\]\s*(\S+)", block)
            link = link_match.group(1).strip() if link_match else ""

            results.append({
                "no": idx,
                "title": title,
                "quotes": quotes,
                "link": link
            })

        return results
    
    def ref_to_html_formatted_2(ref_items):
        """
        parse_reference_txt 결과를
        HTML <p> + <a> 형태로 변환
        """
        html_parts = []

        for item in ref_items:
            html_parts.append(
                f"<p><strong>[{item['no']}] {item['title']}</strong></p>"
            )

            for q in item["quotes"]:
                html_parts.append(
                    f"<p>&nbsp;&nbsp;• {q}</p>"
                )

            if item["link"]:
                html_parts.append(
                    f"<p>&nbsp;&nbsp;🔗 <a href='{item['link']}' target='_blank'>{item['link']}</a></p>"
                )

            html_parts.append("<p></br></p>")

        return "\n".join(html_parts)


    # 참고자료 파일 다시 파싱
    with open(ref_file_path, "r", encoding="utf-8") as f:
        ref_all = parse_reference_txt(f.read())


    def make_reference_p(text: str):
        NS_HP = "{%s}" % NS['hp']

        p = ET.Element(f"{NS_HP}p")
        p.set("styleIDRef", "0")
        p.set("paraPrIDRef", "0")
        p.set("pageBreak", "0")
        p.set("columnBreak", "0")
        p.set("merged", "0")

        run = ET.SubElement(p, f"{NS_HP}run")
        run.set("charPrIDRef", "0")  # 기존 본문과 동일 스타일

        t = ET.SubElement(run, f"{NS_HP}t")
        t.text = text

        return p


    def insert_references_at_end(root, ref_items):
        """
        [참고문헌] 문단을 찾아
        그 아래에 참고자료를 번호 목록 형태로 삽입
        """

        # ① [참고문헌] 문단 찾기
        target_p = None
        p_list = list(root)

        for p in p_list:
            t = p.find(".//hp:t", NS)
            if t is not None and "[참고문헌]" in t.text:
                target_p = p
                break

        if target_p is None:
            print("⚠ [참고문헌] 위치를 찾을 수 없습니다.")
            return

        insert_pos = p_list.index(target_p)

        # ② 참고문헌 문단 삽입
        for i, item in enumerate(ref_items, start=1):
            title = item.get("title", "").strip()
            link = item.get("link", "").strip()

            if link:
                line = f"[{i}] {title}, {link}"
            else:
                line = f"[{i}] {title}"

            new_p = make_reference_p(line)
            root.insert(insert_pos, new_p)
            insert_pos += 1


    insert_references_at_end(root, ref_all)

    # 저장 & 다시 압축
    save_xml(tree)
    zip_hwpx(f"output/{title}.hwpx") 

    print("✔ 제안서 생성 + HWPX 반영 완료:", f"output/{title}.hwpx")


    # ------------------------------------------------------------------------
    # 3. HTML 형태로 파일 생성 (GPT 재호출 없이, 위 텍스트 사용)
    # ------------------------------------------------------------------------

    # HTML 형태로 묶어서 반환
    result_html = f"""
        <h1>사업 제목 : {title}</h1>

        <p>- 키워드 : {keywords_string}</p>
        <p></br></p>
        
        <h3>□ 사업 개요</h3>
        <p>ㅇ 사업 목적 : {purpose}</br></p>
        <p>ㅇ 사업 내용 : {content}</br></p>
        <p></br></p>
     
        <h3>□ 필요성</h3>
        <p>ㅇ {need_parts["a1"]}</br></p>
        <p>&nbsp;&nbsp;- {need_parts["a2"]}</br></p>
        <p>&nbsp;&nbsp;&nbsp;&nbsp;* {need_parts["a3"]}</br></p>
        <p></br></p>
        <p>ㅇ {need_parts["b1"]}</br></p>
        <p>&nbsp;&nbsp;- {need_parts["b2"]}</br></p>
        <p>&nbsp;&nbsp;&nbsp;&nbsp;* {need_parts["b3"]}</br></p>
        <p></br></p>
        <p>ㅇ {need_parts["c1"]}</br></p>
        <p>&nbsp;&nbsp;- {need_parts["c2"]}</br></p>
        <p>&nbsp;&nbsp;&nbsp;&nbsp;* {need_parts["c3"]}</br></p>
        <p></br></p>
        
        <h3>□ 제안 내용</h3>
        <p>{suggestion} + "{title}" + 사업에 국비 반영(국비 300억원)을 요청드림</p>
        <p></br></p>
        <p></br></p>

        <h2>참고</h2>
        <p></br></p>

        <h3>2. 추진 배경 및 필요성</h3>
        <p>{ref_to_html_formatted(ref2)}</p>
        <p></br></p>

        <h3>3. 사업 내용</h3>
        <p>{ref_to_html_formatted(ref3)}</p>
        <p></br></p>

        <h3>4. 기대 효과</h3>
        <p>{ref_to_html_formatted(ref4)}</p>
        <p></br></p>

        <h3><참고 자료></h3>
        <p>{ref_to_html_formatted_2(ref_all)}</p>
        <p></br></p>

     """


    # 최종 결과 HTML 반환
    return result_html



